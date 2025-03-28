<<<<<<< HEAD
from flask import Flask, request, jsonify
from flask_cors import CORS
import yfinance as yf
import finnhub
from duckduckgo_search import DDGS
import json
import datetime
import google.generativeai as genai

app = Flask(__name__)
cors_options = {
    "origins": ["http://localhost:5000"],
    "supports_credentials": True,
    "methods": ["GET", "POST", "PUT", "DELETE", "OPTIONS"],
    "allow_headers": ["Content-Type", "Authorization"]
}
CORS(app, resources={r"/*": cors_options})

genai.configure(api_key='AIzaSyDnGZHEEEZj7m0dGNey9TqGJFtMpN7tmgg')

def gemini_analyze(prompt):
    model = genai.GenerativeModel('gemini-pro')
    response = model.generate_content(prompt)
    return response.text if response else "Analysis not available."

# Agent 1: DataCollectorAgent
class DataCollectorAgent:
    def __init__(self, ticker, finnhub_api_key):
        self.ticker = ticker.upper()
        self.finnhub_client = finnhub.Client(api_key=finnhub_api_key)
    
    def collect_data(self):
        # Collect data using yfinance
        stock = yf.Ticker(self.ticker)
        income_statement = stock.financials
        balance_sheet = stock.balance_sheet
        cash_flow = stock.cashflow

        # Convert financial data to dictionaries, handling None cases
        income_statement_data = income_statement.to_dict() if income_statement is not None else {}
        balance_sheet_data = balance_sheet.to_dict() if balance_sheet is not None else {}
        cash_flow_data = cash_flow.to_dict() if cash_flow is not None else {}

        # Collect company profile using Finnhub
        company_profile = self.finnhub_client.company_profile2(symbol=self.ticker)
        
        today = datetime.datetime.now().strftime('%Y-%m-%d')
        one_year_ago = (datetime.datetime.now() - datetime.timedelta(days=365)).strftime('%Y-%m-%d')
        news = self.finnhub_client.company_news(symbol=self.ticker, _from=one_year_ago, to=today)

        # Collect market news using DuckDuckGo
        ddg_news = []
        with DDGS() as ddgs:
            for result in ddgs.news(f"{self.ticker} market news", max_results=5):
                ddg_news.append(result)

        # Combine all collected data
        collected_data = {
            'income_statement': income_statement_data,
            'balance_sheet': balance_sheet_data,
            'cash_flow': cash_flow_data,
            'company_profile': company_profile,
            'finnhub_news': news,
            'market_news': ddg_news
        }
        print(collected_data)

        return collected_data

# Agent 2: DataAnalyzerAgent
class DataAnalyzerAgent:
    def __init__(self, data, ticker):
        self.data = data
        self.ticker = ticker

    def analyze_data(self):
        analysis = {}

        # 1. Executive Summary
        new_summary = gemini_analyze(f"Generate an executive summary for {self.ticker} based on financial data and market trends.")

        # 2. Income Statement (Profit & Loss)
        analysis['Income Statement'] = self.analyze_income_statement()

        # 3. Balance Sheet
        analysis['Balance Sheet'] = self.analyze_balance_sheet()

        # 4. Cash Flow Statement
        analysis['Cash Flow Statement'] = self.analyze_cash_flow()

        # 5. Key Financial Ratios
        analysis['Key Financial Ratios'] = self.calculate_financial_ratios()

        # 6. Revenue & Cost Breakdown
        analysis['Revenue & Cost Breakdown'] = self.get_revenue_cost_breakdown()

        # 7. Market & Competitive Analysis
        analysis['Market & Competitive Analysis'] = self.market_competitive_analysis()

        # 8. Debt & Capital Structure
        analysis['Debt & Capital Structure'] = self.analyze_debt_capital_structure()

        # 9. Financial Forecasts
        analysis['Financial Forecasts'] = self.financial_forecasts()

        # 10. Risks & Challenges
        analysis['Risks & Challenges'] = self.risks_and_challenges()

        # 11. Management Strategy & Outlook
        analysis['Management Strategy & Outlook'] = self.management_strategy_outlook()
        print(new_summary)

        return analysis

    def generate_executive_summary(self):
        company_name = self.data['company_profile'].get('name', 'the company')

        prompt = f"Generate an executive summary for {company_name} based on the following financial data, balance sheet, cash flow, and market trends: {json.dumps(self.data, indent=4)}"

        summary = gemini_analyze(prompt)

        summary_intro = f"{company_name} has shown significant performance over the past fiscal year."
        summary_str = summary_intro + '\n' + summary
        print(summary_str)
        return summary_str

    def analyze_income_statement(self):
        income_statement = self.data['income_statement']
        if not income_statement:
            return "Income statement data is not available."
        else:
            # Perform analysis (placeholder)
            analysis = f"Analyzed income statement for {self.ticker}."
            return analysis

    def analyze_balance_sheet(self):
        balance_sheet = self.data['balance_sheet']
        if not balance_sheet:
            return "Balance sheet data is not available."
        else:
            # Perform analysis (placeholder)
            analysis = f"Analyzed balance sheet for {self.ticker}."
            return analysis

    def analyze_cash_flow(self):
        cash_flow = self.data['cash_flow']
        if not cash_flow:
            return "Cash flow statement data is not available."
        else:
            # Perform analysis (placeholder)
            analysis = f"Analyzed cash flow statement for {self.ticker}."
            return analysis

    def calculate_financial_ratios(self):
        ratios = {}
        income_statement = self.data['income_statement']
        balance_sheet = self.data['balance_sheet']

        try:
            # Extract necessary financial data
            net_income = income_statement['Net Income'][list(income_statement['Net Income'].keys())[0]]
            total_revenue = income_statement['Total Revenue'][list(income_statement['Total Revenue'].keys())[0]]
            current_assets = balance_sheet['Total Current Assets'][list(balance_sheet['Total Current Assets'].keys())[0]]
            current_liabilities = balance_sheet['Total Current Liabilities'][list(balance_sheet['Total Current Liabilities'].keys())[0]]
            total_debt = balance_sheet['Long Term Debt'][list(balance_sheet['Long Term Debt'].keys())[0]]
            total_equity = balance_sheet['Total Stockholder Equity'][list(balance_sheet['Total Stockholder Equity'].keys())[0]]

            # Profitability ratios
            ratios['Net Margin'] = net_income / total_revenue

            # Liquidity ratios
            ratios['Current Ratio'] = current_assets / current_liabilities

            # Debt ratios
            ratios['Debt to Equity'] = total_debt / total_equity

        except Exception as e:
            ratios['error'] = f"Error calculating ratios: {e}"

        return ratios

    def get_revenue_cost_breakdown(self):
        # Placeholder for actual breakdown
        breakdown = {
            'Revenue Sources': "Details of revenue sources.",
            'Costs': "Details of costs."
        }
        return breakdown

    def market_competitive_analysis(self):
        # Placeholder for actual analysis
        analysis = "Market position, SWOT analysis, and growth potential."
        return analysis

    def analyze_debt_capital_structure(self):
        # Placeholder for actual analysis
        analysis = "Debt levels, financing mix, and repayment ability."
        return analysis

    def financial_forecasts(self):
        # Placeholder for actual forecasts
        forecasts = "Expected revenue and profit growth, risk analysis."
        return forecasts

    def risks_and_challenges(self):
        # Placeholder for actual risks
        risks = "Market risks, legal issues, and operational threats."
        return risks

    def management_strategy_outlook(self):
        # Placeholder for actual strategy
        strategy = "Future plans and business strategy."
        return strategy

# Agent 3: ReportGeneratorAgent
class ReportGeneratorAgent:
    def __init__(self, analysis, ticker):
        self.analysis = analysis
        self.ticker = ticker

    def generate_report(self):
        # Generate JSON report
        report = json.dumps(self.analysis, indent=4)
        return report

@app.route('/generate_report', methods=['POST'])
def generate_report():
    data = request.json
    ticker = data.get('ticker').upper()
    finnhub_api_key = data.get('finnhub_api_key')

    # Create instances of each agent
    data_collector = DataCollectorAgent(ticker, finnhub_api_key)

    # Agent 1 collects data
    collected_data = data_collector.collect_data()
    # Agent 2 analyzes data
    data_analyzer = DataAnalyzerAgent(collected_data, ticker)
    analysis = data_analyzer.analyze_data()

    # Agent 3 generates report
    report_generator = ReportGeneratorAgent(analysis, ticker)
    report = report_generator.generate_report()

    return jsonify(report)

if __name__ == "__main__":
    app.run(debug=True)
=======
import streamlit as st
import base64
import pptx
from pptx.util import Inches, Pt
import os
import time
from dotenv import load_dotenv
import google.generativeai as genai

# Load environment variables
load_dotenv()

# Configure the Generative AI client using your API key
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

# Custom formatting options for the presentation
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

def generate_slide_titles(topic):
    """Generates slide titles for a given topic using the Generative AI client."""
    prompt = f"Generate 5 slide titles for the topic '{topic}'."
    response = genai.generate_text(prompt=prompt)
    titles = response.result.split("\n")
    return [title.strip() for title in titles if title.strip()]

def generate_slide_content(slide_title):
    """Generates detailed slide content based on a slide title."""
    prompt = f"Generate detailed content for a slide titled '{slide_title}'."
    response = genai.generate_text(prompt=prompt)
    return response.result.strip()

def create_presentation(topic, slide_titles, slide_contents):
    """Creates a PowerPoint presentation from slide titles and their content."""
    prs = pptx.Presentation()
    # Create a title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    # Create a slide for each title-content pair
    for title, content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.shapes.placeholders[1].text = content

        # Set custom font sizes
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE

    # Save the presentation in a directory
    output_dir = "generated_ppt"
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    filename = os.path.join(output_dir, f"{topic}_presentation.pptx")
    prs.save(filename)
    return filename

def get_ppt_download_link(filepath):
    """Generates an HTML download link for the PPT file."""
    with open(filepath, "rb") as f:
        ppt_data = f.read()
    b64 = base64.b64encode(ppt_data).decode()
    return (
        f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" '
        f'download="{os.path.basename(filepath)}">Download Presentation</a>'
    )

def main():
    st.title("PowerPoint Presentation Generator using Google Gen AI (Gemini)")
    topic = st.text_input("Enter the topic for your presentation:")
    if st.button("Generate Presentation") and topic:
        st.info("Generating presentation, please wait...")
        start_time = time.time()

        slide_titles = generate_slide_titles(topic)
        slide_contents = [generate_slide_content(title) for title in slide_titles]
        ppt_file = create_presentation(topic, slide_titles, slide_contents)
        elapsed = time.time() - start_time
        st.success(f"Presentation generated in {elapsed:.2f} seconds!")
        st.markdown(get_ppt_download_link(ppt_file), unsafe_allow_html=True)

if _name_ == "_main_":
    main()
>>>>>>> 12b032c14be6fa1689dfed0a2e759f53fbb400de
